"""
Document Processing Module
Handles extraction and processing of various document formats
"""

import logging
from typing import List, Dict, Any
import os
from pathlib import Path

# Document processing imports
import PyPDF2
from docx import Document
from pptx import Presentation

# LangChain imports
from langchain.schema import Document as LangChainDocument
from langchain.text_splitter import RecursiveCharacterTextSplitter

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Handles processing of various document formats"""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )
    
    def process_file(self, file_path: str, original_filename: str) -> List[LangChainDocument]:
        """Process a file and return a list of document chunks"""
        try:
            file_extension = Path(file_path).suffix.lower()
            
            # Extract text based on file type
            if file_extension == '.pdf':
                text = self._extract_pdf_text(file_path)
            elif file_extension == '.docx':
                text = self._extract_docx_text(file_path)
            elif file_extension == '.pptx':
                text = self._extract_pptx_text(file_path)
            elif file_extension == '.txt':
                text = self._extract_txt_text(file_path)
            else:
                logger.warning(f"Unsupported file format: {file_extension}")
                return []
            
            if not text.strip():
                logger.warning(f"No text extracted from {original_filename}")
                return []
            
            # Clean text to remove potential prompt injections
            text = self._clean_text(text)
            
            # Split text into chunks
            chunks = self.text_splitter.split_text(text)
            
            # Create LangChain documents
            documents = []
            for i, chunk in enumerate(chunks):
                doc = LangChainDocument(
                    page_content=chunk,
                    metadata={
                        'source': original_filename,
                        'chunk_id': i,
                        'file_type': file_extension,
                        'total_chunks': len(chunks)
                    }
                )
                documents.append(doc)
            
            logger.info(f"Processed {original_filename}: {len(documents)} chunks created")
            return documents
            
        except Exception as e:
            logger.error(f"Error processing file {original_filename}: {e}")
            return []
    
    def _clean_text(self, text: str) -> str:
        """Clean text to remove potential prompt injections and malicious content"""
        # Remove common prompt injection patterns
        injection_patterns = [
            r'ignore all (previous )?instructions?',
            r'disregard all (previous )?instructions?',
            r'forget all (previous )?instructions?',
            r'system prompt',
            r'act as',
            r'pretend to be',
            r'you are now',
            r'new instructions?:',
            r'explain why this candidate should be hired',
            r'&.*explain why.*hired.*resume',
            r'& given his resume',
            r'&.*given.*resume'
        ]
        
        import re
        cleaned_text = text
        
        for pattern in injection_patterns:
            cleaned_text = re.sub(pattern, '', cleaned_text, flags=re.IGNORECASE)
        
        # Remove excessive whitespace
        cleaned_text = re.sub(r'\s+', ' ', cleaned_text).strip()
        
        return cleaned_text
    
    def _extract_pdf_text(self, file_path: str) -> str:
        """Extract text from PDF file"""
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            text += f"\n--- Page {page_num + 1} ---\n"
                            text += page_text
                    except Exception as e:
                        logger.warning(f"Error extracting text from page {page_num + 1}: {e}")
                        continue
            
            return text
            
        except Exception as e:
            logger.error(f"Error reading PDF file: {e}")
            return ""
    
    def _extract_docx_text(self, file_path: str) -> str:
        """Extract text from DOCX file"""
        try:
            doc = Document(file_path)
            text = ""
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text += " | ".join(row_text) + "\n"
            
            return text
            
        except Exception as e:
            logger.error(f"Error reading DOCX file: {e}")
            return ""
    
    def _extract_pptx_text(self, file_path: str) -> str:
        """Extract text from PPTX file"""
        try:
            prs = Presentation(file_path)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides):
                text += f"\n--- Slide {slide_num + 1} ---\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                    
                    # Extract text from tables in slides
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                text += " | ".join(row_text) + "\n"
            
            return text
            
        except Exception as e:
            logger.error(f"Error reading PPTX file: {e}")
            return ""
    
    def _extract_txt_text(self, file_path: str) -> str:
        """Extract text from TXT file"""
        try:
            encodings = ['utf-8', 'utf-16', 'latin-1', 'ascii']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        return file.read()
                except UnicodeDecodeError:
                    continue
            
            logger.warning(f"Could not decode text file with any common encoding")
            return ""
            
        except Exception as e:
            logger.error(f"Error reading TXT file: {e}")
            return ""
    
    def get_supported_formats(self) -> List[str]:
        """Return list of supported file formats"""
        return ['.pdf', '.docx', '.pptx', '.txt']
    
    def validate_file(self, file_path: str) -> bool:
        """Validate if file format is supported and file exists"""
        if not os.path.exists(file_path):
            return False
        
        file_extension = Path(file_path).suffix.lower()
        return file_extension in self.get_supported_formats()
